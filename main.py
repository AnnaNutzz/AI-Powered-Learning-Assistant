# Flask web framework and utility imports
from flask import Flask, render_template, request, redirect, url_for, flash, session, jsonify
import sqlite3  # SQLite database
import os  # File system and environment variables
from werkzeug.utils import secure_filename  # Secure file handling
from werkzeug.security import generate_password_hash, check_password_hash  # Password hashing
import pdfplumber  # For PDF text extraction
from pptx import Presentation  # For PPTX text extraction
from dotenv import load_dotenv  # Load environment variables from .env
from notion_client import Client  # Notion API
from twilio.rest import Client as TwilioClient  # Twilio for SMS
import random
import torch
from transformers import pipeline  # Hugging Face summarizer
import ollama  # Chat-based summarization with ollama

# Load environment variables from .env
load_dotenv()

# Use GPU if available, otherwise CPU
summarizer = pipeline("summarization", model="facebook/bart-large-cnn", device=0 if torch.cuda.is_available() else -1)

# Initialize Flask app and configuration
app = Flask(__name__)
app.secret_key = os.getenv("FLASK_SECRET_KEY", os.urandom(24).hex())
app.config['UPLOAD_FOLDER'] = 'uploads/'
app.config['ALLOWED_EXTENSIONS'] = {'pdf', 'pptx', 'txt'}

# Initialize external service credentials
NOTION_TOKEN = os.getenv("NOTION_TOKEN")
NOTION_DATABASE_ID = os.getenv("NOTION_DATABASE_ID")
TWILIO_PHONE_NUMBER = os.getenv("TWILIO_PHONE_NUMBER")
TWILIO_SID = os.getenv("TWILIO_SID")
TWILIO_AUTH_TOKEN = os.getenv("TWILIO_AUTH_TOKEN")

# Initialize Twilio and Notion clients
twilio_client = TwilioClient(TWILIO_SID, TWILIO_AUTH_TOKEN)
notion = Client(auth=os.getenv("NOTION_SECRET"))

# Initialize SQLite database with user and revision tables
def init_db():
    with sqlite3.connect("users.db") as conn:
        cursor = conn.cursor()
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS users (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                username TEXT UNIQUE NOT NULL,
                password TEXT NOT NULL,
                parent_phone TEXT NOT NULL,
                learning_type TEXT NOT NULL,
                learning_speed TEXT NOT NULL
            )
        ''')
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS revision (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                username TEXT NOT NULL,
                content TEXT NOT NULL,
                file_path TEXT
            )
        ''')
        conn.commit()


# Validate file extension
def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

# Extract text from PDF
def extract_pdf_text(file_path):
    try:
        with pdfplumber.open(file_path) as pdf:
            text = "".join(page.extract_text() or "" for page in pdf.pages)
        return text
    except Exception as e:
        return f"Error extracting PDF text: {str(e)}"

# Extract text from PPTX
def extract_pptx_text(file_path):
    try:
        prs = Presentation(file_path)
        text = "".join(shape.text + "\n" for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
        return text
    except Exception as e:
        return f"Error extracting PPTX text: {str(e)}"

import ollama

import ollama

ollama_client = ollama.Client()

def generate_quick_notes(text, learning_speed):
    try:
        input_length = len(text.split())
        if input_length < 10:
            return "Text too short to summarize."

        max_length = min(int(input_length * 0.8), 150)
        min_length = int(max_length * 0.3)

        if learning_speed == "Slow":
            max_length = min(max_length, 50)
        elif learning_speed == "Average":
            max_length = min(max_length, 100)

        # Include specific summarization prompt for better results
        prompt = f"Summarize the following text in {min_length}-{max_length} words:\n\n{text}"

        response = ollama_client.chat(
            model="gemma:2b",
            messages=[{"role": "user", "content": prompt}]
        )

        return response['message']['content']
    
    except Exception as e:
        return f"Error generating notes: {str(e)}"

# Work in progress
# Add notes to Notion
def add_notes_to_notion(notes, username):
    try:
        notion.pages.create(
            parent={"database_id": NOTION_DATABASE_ID},
            properties={
                "Name": {"title": [{"text": {"content": f"Notes for {username}"}}]},
                "Notes": {"rich_text": [{"text": {"content": notes}}]}
            }
        )
        return True
    except Exception as e:
        return f"Error adding to Notion: {str(e)}"

# Routes
@app.route("/")
def home():
    if 'username' not in session:
        return redirect(url_for('login'))
    with sqlite3.connect("users.db") as conn:
        cursor = conn.cursor()
        cursor.execute("SELECT learning_type FROM users WHERE username = ?", (session['username'],))
        learning_type = cursor.fetchone()[0]
    if not learning_type or request.args.get("redo", "false").lower() == "true":
        return redirect(url_for("quiz"))
    return render_template("home.html", learning_type=learning_type)

@app.route('/dashboard')
def dashboard():
    if 'username' not in session:
        return redirect(url_for('login'))
    username = session['username']
    with sqlite3.connect("users.db") as conn:
        cursor = conn.cursor()
        cursor.execute("SELECT learning_type, learning_speed FROM users WHERE username = ?", (username,))
        result = cursor.fetchone()
    learning_type = result[0] if result else "Not determined"
    learning_speed = result[1] if result else "Not determined"
    suggestions = {
        "Watching": "Watch videos or tutorials.",
        "Reading": "Read articles and textbooks.",
        "Listening": "Listen to podcasts or lectures.",
        "Doing": "Try hands-on exercises."
    }.get(learning_type, "Complete the quiz to get suggestions.")
    return render_template('dashboard.html', learning_type=learning_type, learning_speed=learning_speed, suggestions=[suggestions])

@app.route('/register', methods=['GET', 'POST'])
def register():
    if request.method == 'POST':
        username = request.form.get('username')
        password = request.form.get('password')
        parent_phone = request.form.get('parent_phone')
        if not all([username, password, parent_phone]):
            flash('All fields are required.', 'danger')
            return render_template('register.html')
        if len(password) < 8:
            flash('Password must be at least 8 characters long.', 'danger')
            return render_template('register.html')
        hashed_password = generate_password_hash(password)
        try:
            with sqlite3.connect("users.db") as conn:
                cursor = conn.cursor()
                cursor.execute("INSERT INTO users (username, password, parent_phone, learning_type, learning_speed) VALUES (?, ?, ?, ?, ?)",
                              (username, hashed_password, parent_phone, '', 'Average'))
                conn.commit()
                flash('Registration successful! Please login.', 'success')
                return redirect(url_for('login'))
        except sqlite3.IntegrityError:
            flash('Username already exists.', 'danger')
        except Exception as e:
            flash(f'Error: {str(e)}', 'danger')
    return render_template('register.html')

@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        username = request.form.get('username')
        password = request.form.get('password')
        if not all([username, password]):
            flash('Username and password are required.', 'danger')
            return render_template('login.html')
        with sqlite3.connect("users.db") as conn:
            cursor = conn.cursor()
            cursor.execute("SELECT password, parent_phone, learning_speed FROM users WHERE username = ?", (username,))
            result = cursor.fetchone()
            if result and check_password_hash(result[0], password):
                session['username'] = username
                session['parent_phone'] = result[1]
                session['learning_speed'] = result[2] if len(result) > 2 else "Average"
                flash('Login successful!', 'success')
                return redirect(url_for('home'))
            flash('Invalid username or password.', 'danger')
    return render_template('login.html')

@app.route('/logout')
def logout():
    session.clear()
    flash('Logged out successfully.', 'success')
    return redirect(url_for('login'))

@app.route('/send_sms', methods=['POST'])
def send_sms():
    data = request.get_json()
    phone_number = data.get('phone_number')
    message = data.get('message')
    if not all([phone_number, message]):
        return jsonify({"error": "Phone number and message are required."}), 400
    try:
        twilio_client.messages.create(
            body=message,
            from_=TWILIO_PHONE_NUMBER,
            to=phone_number
        )
        return jsonify({"message": "SMS sent successfully!"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/select_grade_level', methods=['GET', 'POST'])
def select_grade_level():
    if 'username' not in session:
        return redirect(url_for('login'))
    if request.method == 'POST':
        grade_level = request.form.get('grade_level')
        if not grade_level:
            flash('Please select a grade level.', 'danger')
            return render_template('select_grade_level.html')
        session['grade_level'] = grade_level
        return redirect(url_for('quiz'))
    return render_template('select_grade_level.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'username' not in session:
        return jsonify({"error": "Please log in."}), 401
    if 'file' not in request.files:
        return jsonify({"error": "No file part."}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file."}), 400
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        user_folder = os.path.join(app.config['UPLOAD_FOLDER'], session['username'])
        os.makedirs(user_folder, exist_ok=True)
        file_path = os.path.join(user_folder, filename)
        file.save(file_path)
        # Process file
        content = ""
        if filename.endswith('.pdf'):
            content = extract_pdf_text(file_path)
        elif filename.endswith(('.ppt', '.pptx')):
            content = extract_pptx_text(file_path)
        elif filename.endswith('.txt'):
            with open(file_path, 'r') as f:
                content = f.read()
        # Generate quick notes
        learning_speed = session.get('learning_speed', "Average")
        try:
            summary = generate_quick_notes(content, learning_speed) if content else "No content to summarize."
        except Exception as e:
            summary = f"Error generating notes: {str(e)}"


        # Save to revision table
        with sqlite3.connect("users.db") as conn:
            cursor = conn.cursor()
            cursor.execute("INSERT INTO revision (username, content, file_path) VALUES (?, ?, ?)",
                          (session['username'], summary, file_path))
            conn.commit()
        # Add to Notion
        notion_result = add_notes_to_notion(summary, session['username'])
        if "Error" in str(notion_result):
            flash(notion_result, 'warning')
        return jsonify({"message": "File uploaded and processed successfully!", "summary": summary})
    return jsonify({"error": "Invalid file type."}), 400

@app.route('/quiz', methods=['GET', 'POST'])
def quiz():
    if 'username' not in session:
        return redirect(url_for('login'))
    if not session.get('grade_level'):
        return redirect(url_for('select_grade_level'))

    questions = [
        {
            "question": "How do you prefer to learn new material?",
            "options": ["Reading books or articles", "Watching videos", "Listening to lectures or podcasts", "Hands-on practice"],
            "type": ["Reading", "Watching", "Listening", "Doing"]
        },
        {
            "question": "What helps you retain information best?",
            "options": ["Taking notes while reading", "Watching demonstrations", "Discussing with others", "Doing exercises"],
            "type": ["Reading", "Watching", "Listening", "Doing"]
        },
        {
            "question": "How do you prepare for exams?",
            "options": ["Reading textbooks", "Watching review videos", "Listening to recordings", "Practicing problems"],
            "type": ["Reading", "Watching", "Listening", "Doing"]
        }
    ]

    if request.method == 'POST':
        answers = request.form
        counts = {"Reading": 0, "Watching": 0, "Listening": 0, "Doing": 0}
        for i, q in enumerate(questions):
            selected = answers.get(f'q{i}')
            if selected in q['type']:
                counts[selected] += 1
        # Determine dominant learning type
        learning_type = max(counts, key=counts.get)
        learning_speed = random.choice(["Slow", "Average", "Fast"])  # You can update this to be quiz-based

        # Save to DB
        with sqlite3.connect("users.db") as conn:
            cursor = conn.cursor()
            cursor.execute("UPDATE users SET learning_type = ?, learning_speed = ? WHERE username = ?",
                           (learning_type, learning_speed, session['username']))
            conn.commit()
        session['learning_speed'] = learning_speed
        flash("Quiz completed! Personalized learning plan is ready.", "success")
        return redirect(url_for('dashboard'))

    return render_template('quiz.html', questions=questions)

@app.route('/quiz_result', methods=['POST'])
def quiz_result():
    if 'username' not in session:
        return redirect(url_for('login'))
    questions = session.get('questions', [])
    if not questions:
        flash('No quiz questions found.', 'danger')
        return redirect(url_for('quiz'))
    # Count learning style preferences
    style_count = {"Reading": 0, "Watching": 0, "Listening": 0, "Doing": 0}
    for i, question in enumerate(questions):
        answer = request.form.get(f'answer{i}')
        if not answer or answer not in question['options']:
            flash('Invalid or missing answer.', 'danger')
            return redirect(url_for('quiz'))
        style = question['type'][question['options'].index(answer)]
        style_count[style] += 1
    # Determine learning type
    learning_type = max(style_count, key=style_count.get)
    # Self-reported learning speed
    learning_speed = request.form.get('learning_speed', 'Average')
    if learning_speed not in ['Slow', 'Average', 'Fast']:
        learning_speed = 'Average'
    # Update database
    try:
        with sqlite3.connect("users.db") as conn:
            cursor = conn.cursor()
            cursor.execute('''UPDATE users SET learning_type = ?, learning_speed = ? WHERE username = ?''',
                          (learning_type, learning_speed, session['username']))
            conn.commit()
    except Exception as e:
        flash(f"Error updating profile: {str(e)}", 'danger')
        return redirect(url_for('quiz'))
    # Send SMS to parent
    try:
        message = f"Your child {session['username']} completed the quiz. Learning type: {learning_type}, Learning speed: {learning_speed}."
        twilio_client.messages.create(
            body=message,
            from_=TWILIO_PHONE_NUMBER,
            to=session['parent_phone']
        )
    except Exception as e:
        flash(f"Error sending SMS: {str(e)}", 'warning')
    flash('Quiz completed successfully!', 'success')
    return redirect(url_for('dashboard'))


from transformers import pipeline

# Initialize summarizer pipeline
summarizer = pipeline("summarization", model="facebook/bart-large-cnn")

# Function to chunk content if it's too large
def chunk_text(text, chunk_size=1024):
    # Split the text into manageable chunks
    return [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]

# Function to summarize large files
def summarize_text(text):
    input_length = len(text.split())  # Count the words in the text
    
    # Set max_length to a reasonable value based on input length, but not too short
    max_length = min(input_length, 150)  # Set a max_length limit based on the input length
    
    # Ensure min_length is smaller than max_length
    min_length = max(30, int(max_length * 0.2))  # Set min_length to be at least 20% of max_length
    
    try:
        # Summarize the text
        summary = summarizer(text, max_length=max_length, min_length=min_length, do_sample=False)
        
        # Extract the summary text from the list of summaries
        summary_text = summary[0]['summary_text']  # Assuming the summary is in the first element of the list
        
        return summary_text.strip()  # Return the summary text after stripping any leading/trailing whitespace
    except Exception as e:
        print(f"Error generating notes: {e}")
        return None

# Example usage
uploaded_text = "Your long file content here..."  # Use the uploaded file content here
summary = summarize_text(uploaded_text)

@app.route('/revise', methods=['GET'])
def revise():
    if 'username' not in session:
        return redirect(url_for('login'))
    with sqlite3.connect("users.db") as conn:
        cursor = conn.cursor()
        cursor.execute("SELECT content, file_path FROM revision WHERE username = ?", (session['username'],))
        revisions = cursor.fetchall()
    if revisions:
        random_revision = random.choice(revisions)
        return render_template('revision.html', file={'notes': random_revision[0], 'filename': random_revision[1]})
    return jsonify({'message': 'No revision files available'}), 400

@app.route('/send_revision_done_message', methods=['POST'])
def send_revision_done_message():
    if 'username' not in session or 'parent_phone' not in session:
        return jsonify({'message': 'Please log in.'}), 401
    try:
        message = f"Your child {session['username']} has completed their revision."
        twilio_client.messages.create(
            body=message,
            from_=TWILIO_PHONE_NUMBER,
            to=session['parent_phone']
        )
    except Exception as e:
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    init_db()
    app.run(debug=True)
